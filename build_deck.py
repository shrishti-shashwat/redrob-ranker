"""Fill the mandatory India.Runs template with the redrob-ranker approach.

Keeps the template untouched except for text content: questions become
bold lead-ins, answers go beneath in the same Manrope/purple styling.
Slide 7 gets a native-shapes architecture diagram.

    python build_deck.py            # -> deck_submission.pptx
"""

import copy

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
INK, PURPLE = "202729", "7D45E0"

TEAM_NAME = "Starva"
TEAM_LEADER = "Shrishti"
GITHUB_URL = "<<GITHUB REPO URL>>"
SANDBOX_URL = "<<SANDBOX LINK>>"


def el(tag, **attrs):
    e = etree.SubElement
    node = etree.Element(f"{{{A}}}{tag}")
    for k, v in attrs.items():
        node.set(k, str(v))
    return node


def make_para(text, *, bold=False, size=11.0, bullet="dot", indent_in=0.0,
              color=INK, spacing=115):
    """Build an <a:p> matching the template's Manrope styling."""
    p = el("p")
    pPr = el("pPr", lvl=0, rtl=0, algn="l",
             marL=int(Inches(0.32 + indent_in)), indent=-Inches(0.22))
    ln = el("lnSpc"); ln.append(el("spcPct", val=spacing * 1000)); pPr.append(ln)
    for t, v in (("spcBef", 600 if bullet == "dot" else 0), ("spcAft", 0)):
        sp = el(t); sp.append(el("spcPts", val=v)); sp.append
        pPr.append(sp)
    if bullet == "dot":
        bc = el("buClr"); bc.append(el("srgbClr", val=PURPLE)); pPr.append(bc)
        pPr.append(el("buSzPts", val=int(size * 100)))
        pPr.append(el("buFont", typeface="Manrope SemiBold"))
        pPr.append(el("buChar", char="●"))
    elif bullet == "dash":
        bc = el("buClr"); bc.append(el("srgbClr", val=PURPLE)); pPr.append(bc)
        pPr.append(el("buSzPts", val=int(size * 100)))
        pPr.append(el("buFont", typeface="Manrope SemiBold"))
        pPr.append(el("buChar", char="–"))
    else:
        pPr.append(el("buNone"))
    p.append(pPr)

    r = el("r")
    rPr = el("rPr", lang="en-GB", sz=int(size * 100))
    if bold:
        rPr.set("b", "1")
    fill = el("solidFill"); fill.append(el("srgbClr", val=color)); rPr.append(fill)
    for tag in ("latin", "ea", "cs", "sym"):
        rPr.append(el(tag, typeface="Manrope SemiBold"))
    r.append(rPr)
    t = el("t"); t.text = text; r.append(t)
    p.append(r)
    return p


def fill_body(slide, paras):
    shape = body_of(slide)
    tx = shape.text_frame._txBody
    for p in tx.findall(f"{{{A}}}p"):
        tx.remove(p)
    for spec in paras:
        tx.append(make_para(spec[0], **spec[1]))


Q = lambda t: (t, dict(bold=True, size=11.5, bullet="dot"))
ANS = lambda t: (t, dict(size=10.5, bullet="dash", indent_in=0.28))
STEP = lambda t: (t, dict(size=11.0, bullet="dot"))


def body_of(slide, must_contain=None):
    boxes = [sh for sh in slide.shapes if sh.has_text_frame]
    if must_contain:
        boxes = [b for b in boxes if must_contain in b.text_frame.text]
    return max(boxes, key=lambda b: b.width * b.height)


prs = Presentation("deck_template.pptx")
S = prs.slides

# ---- Slide 1: title fields -------------------------------------------------
for sh in S[0].shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    add = {"Team Name :": f" {TEAM_NAME}",
           "Team Leader Name :": f" {TEAM_LEADER}",
           "Problem Statement :": (" Intelligent candidate discovery & ranking"
                                   " — rank 100K profiles for a Senior AI"
                                   " Engineer JD the way a great recruiter"
                                   " would, not by keyword matching.")}.get(t)
    if add:
        para = sh.text_frame.paragraphs[0]
        run = para.runs[0]
        new = copy.deepcopy(run._r)
        new.find(f"{{{A}}}t").text = add
        run._r.addnext(new)

# ---- Slide 2: Solution Overview -------------------------------------------
fill_body(S[1], [
    Q("What is your proposed solution?"),
    ANS("A hybrid ranking engine: deterministic profile-consistency checks → "
        "corroboration-gated structured JD-fit features → semantic (MiniLM "
        "bi-encoder, 5 JD facets) + BM25 fusion → behavioral availability "
        "multiplier → top-100 with fact-grounded reasoning per candidate."),
    ANS("All heavy computation is precomputed offline and shipped as artifacts; "
        "the timed ranking step runs in ~2 s on CPU against a 300 s budget."),
    Q("What differentiates it from traditional candidate matching systems?"),
    ANS("Keywords are treated as claims, not evidence. A listed skill scores only "
        "when the title or career history corroborates it — measured on this "
        "pool, P(PyTorch | ML-titled) = 21.6% vs 1.2% for everyone else — so "
        "keyword stuffers get exactly zero credit."),
    ANS("What people actually built (career narratives), where they built it "
        "(product vs consulting vs AI-native companies) and how they engage "
        "(recruiter signals) drive the rank; impossible profiles are detected "
        "and excluded; every score is explainable with no LLM at runtime."),
])

# ---- Slide 3: JD Understanding ---------------------------------------------
fill_body(S[2], [
    Q("Key requirements extracted from the JD"),
    ANS("Core: built retrieval / ranking / recommendation systems in production; "
        "evaluates ranking quality rigorously (NDCG, A/B); 5–9 yrs treated as "
        "a soft band; product-company shipping culture; LLM fine-tuning a plus."),
    ANS("Logistics: Noida/Pune or willing to relocate (no visa sponsorship); "
        "notice ≤30 days buyable. Disqualifiers are graded like the JD's own "
        "language: consulting-only and research-only careers = hard kills "
        "('will not move forward'); job-hopping, CV/speech-only, LangChain-only "
        "= heavy penalties ('probably not')."),
    Q("Which signals matter most / fit beyond keyword matching?"),
    ANS("Career-history descriptions are the primary evidence, then title "
        "archetypes and company taxonomy; redrob behavioral signals — "
        "last-active, recruiter response rate, open-to-work, notice period, "
        "relocation — modulate the skill-match score multiplicatively, "
        "exactly as the signals doc prescribes."),
])

# ---- Slide 4: Ranking Methodology ------------------------------------------
fill_body(S[3], [
    Q("How does the system retrieve, score and rank?"),
    ANS("No hard retrieval gate — all 100K candidates receive all scores (a "
        "lexical gate would silently drop plain-language strong profiles). "
        "Semantic: MiniLM-L6-v2 embeddings vs 5 authored JD facets, mean of "
        "top-2 similarities. Lexical: BM25 vs a JD-derived query. Structured: "
        "~20 features (title archetype, retrieval evidence, production & "
        "evaluation signals, corroborated skills, product ratio, experience band)."),
    Q("Models, algorithms, heuristics — and why"),
    ANS("Bi-encoder (22M params, CPU-friendly), BM25, and hand-weighted linear "
        "scoring. Interpretable by design: no labelled ground truth exists, so a "
        "trained learning-to-rank model would only learn our assumptions back, "
        "with less transparency."),
    Q("How are signals combined into the final ranking?"),
    ANS("z-normalise → sigmoid → fuse 0.45·structured + 0.40·semantic + "
        "0.15·BM25; multiplicative kills (honeypot ×0.02, consulting-only / "
        "research-only ×0.05, non-tech title ×0.10); × availability — a "
        "weighted geometric mean of recency, response rate, open-to-work, "
        "notice, location, interview completion (floor 0.25, so behavioral "
        "twins actually separate). Deterministic tie-break; every date anchored "
        "to the dataset's frozen clock, 2026-06-01."),
])

# ---- Slide 5: Explainability & Data Validation ------------------------------
fill_body(S[4], [
    Q("How are ranking decisions explained?"),
    ANS("Programmatic generation from extracted facts only: real company names, "
        "matched retrieval terms, corroborated skills, signal values. Concerns "
        "are disclosed systematically — notice >30 days, dormancy, not "
        "open-to-work, overseas without relocation — and tone tracks rank."),
    Q("How do you prevent hallucinations or unsupported justifications?"),
    ANS("No LLM at runtime. Every clause renders from verified profile fields, "
        "so an unsupported claim is structurally impossible — and the output "
        "is deterministic, reproducing byte-for-byte across machines."),
    Q("How are inconsistent, low-quality or suspicious profiles handled?"),
    ANS("68 planted honeypots caught by three record-internal checks: 'expert' "
        "skill with 0 months used; current job implying today ≠ 2026-06; "
        "career span contradicting stated experience by >3 yrs (with the "
        "±1-month stint tolerance the data requires — 64,649 stints carry "
        "ordinary calendar rounding)."),
    ANS("Keyword stuffers die at the corroboration gate; behavioral twins are "
        "separated by the availability multiplier; the pre-submission audit "
        "enforces zero honeypots in the top 100."),
])

# ---- Slide 6: End-to-End Workflow -------------------------------------------
fill_body(S[5], [
    STEP("1.  JD analysis — requirements, graded disqualifiers, and 5 semantic "
         "facet queries authored from the role description"),
    STEP("2.  Offline precompute (~90 min, permitted by spec §10.3) — MiniLM "
         "embeddings for all 100K profiles, BM25 scores, and a structured-"
         "feature / availability / honeypot cache"),
    STEP("3.  Timed ranking step (~2 s of the 300 s budget) — load artifacts, "
         "fuse scores, apply kills and penalties, multiply availability, "
         "take the top 100"),
    STEP("4.  Reasoning — fact-grounded explanation per candidate, concerns "
         "included"),
    STEP("5.  Validation & audit — official validator, zero-honeypot check, "
         "twin-ordering check, manual top-50 review → submission.csv"),
])

# ---- Slide 7: System Architecture (diagram) ---------------------------------
sl = S[6]


def box(x, y, w, h, title, lines, fill="F4F0FB", line=PURPLE, title_color=INK):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    sh.line.color.rgb = RGBColor.from_string(line); sh.line.width = Pt(1.2)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tx = tf._txBody
    for p in tx.findall(f"{{{A}}}p"):
        tx.remove(p)
    tx.append(make_para(title, bold=True, size=10.5, bullet=None,
                        color=title_color, spacing=100))
    for ln in lines:
        tx.append(make_para(ln, size=9, bullet=None, color=INK, spacing=100))
    return sh


def arrow(x1, y1, x2, y2):
    ln = sl.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = RGBColor.from_string(PURPLE); ln.line.width = Pt(1.75)
    ln.line._get_or_add_ln().append(
        etree.fromstring(f'<a:tailEnd xmlns:a="{A}" type="triangle"/>'))


box(0.45, 1.30, 1.95, 1.30, "INPUTS", ["candidates.jsonl", "(100K, 487 MB)",
                                       "job_description"], fill="EDEDF2")
box(0.45, 3.45, 1.95, 1.40, "TRAP DEFENSES", ["68 honeypots flagged",
    "corroboration gate", "twin separation"], fill="FBEFEF", line="B85042")

box(2.85, 1.30, 2.75, 1.85, "OFFLINE — run once (~90 min)",
    ["MiniLM embeddings 100K×384", "BM25 scores vs JD query",
     "features.pkl — fit, availability,", "honeypot flags  (spec §10.3)"])
box(2.85, 3.45, 2.75, 1.40, "JD FACETS",
    ["5 authored 'ideal candidate'", "queries: retrieval, eval rigor,",
     "product shipping, LLM, NLP"])

box(6.05, 1.30, 2.45, 2.20, "RANKING — timed step (~2 s)",
    ["budget: 300 s, CPU, no network", "fuse 0.45 struct + 0.40 sem",
     "+ 0.15 BM25 (z → sigmoid)",
     "× kills: honeypot, consulting/", "research-only, non-tech",
     "× availability (geo-mean, fl 0.25)"], fill="EFE9FC")

box(6.05, 3.85, 2.45, 1.00, "OUTPUT", ["submission.csv — top 100",
    "+ fact-grounded reasoning"], fill="E9F4EE", line="2C5F2D")

arrow(2.40, 1.95, 2.85, 1.95)   # inputs -> offline
arrow(2.40, 3.60, 2.90, 3.20)   # trap defenses -> offline (flags + gate live there)
arrow(5.60, 2.20, 6.05, 2.20)   # offline -> ranking
arrow(4.22, 3.45, 4.22, 3.15)   # facets -> offline (embedded together)
arrow(7.27, 3.50, 7.27, 3.85)   # ranking -> output

# ---- Slide 8: Results & Performance -----------------------------------------
fill_body(S[7], [
    Q("What results demonstrate ranking quality?"),
    ANS("Top-100 is what a recruiter would expect for this JD: search / recsys / "
        "ML engineers, 4–9 yrs, at AI-native startups and product companies "
        "— surfaced from a pool that is 64% non-tech and only ~1% ML-titled."),
    ANS("0 of 68 detectable honeypots in the top 100 (DQ threshold: >10); "
        "behavioral-twin ordering verified; every risky row carries an honest, "
        "data-backed concern; facts in reasonings spot-checked against raw "
        "profiles with zero errors."),
    Q("How are the runtime and compute constraints met?"),
    ANS("Timed ranking step: ~2 s wall-clock vs the 300 s budget, CPU-only, no "
        "network, well under 16 GB. Offline precompute (~90 min) ships as "
        "artifacts per spec §10.3."),
    ANS("All recency math anchors to the dataset's frozen clock (2026-06-01), "
        "never the wall clock — so reproduction is byte-identical on any "
        "machine, any day."),
])

# ---- Slide 9: Technologies Used ----------------------------------------------
fill_body(S[8], [
    ("Python 3.11 + NumPy — vectorised scoring; 100K×384 similarity is a "
     "single matrix multiply", dict(size=11, bullet="dot")),
    ("sentence-transformers (MiniLM-L6-v2, 22M params) — the best "
     "quality-per-CPU-second bi-encoder at this scale; embeds the full pool "
     "offline without a GPU", dict(size=11, bullet="dot")),
    ("Custom BM25 (NumPy, K1=1.5 b=0.75) — document-length-aware lexical "
     "signal fused alongside semantics; no black-box dependency", dict(size=11, bullet="dot")),
    ("Python stdlib for feature extraction and reasoning generation — "
     "deterministic, no-network, zero-hallucination runtime", dict(size=11, bullet="dot")),
    ("NPY / pickle artifact cache — moves every heavy step out of the timed "
     "window (127–311 s of JSONL parsing → ~2 s)", dict(size=11, bullet="dot")),
    ("Deliberately no GPU, no external APIs, no LLM at inference — the "
     "constraints reward small, sharp and reproducible over large and fragile",
     dict(size=11, bullet="dot")),
])

# ---- Slide 10: Submission Assets ----------------------------------------------
fill_body(S[9], [
    (f"GitHub repository: {GITHUB_URL} — full source, README with "
     "one-command reproduction, precompute scripts, 5 commits of real "
     "iteration history", dict(size=11.5, bullet="dot")),
    ("Ranked output: submission.csv — top-100, passes the official "
     "validator", dict(size=11.5, bullet="dot")),
    (f"Sandbox: {SANDBOX_URL} — HuggingFace Space running the ranker "
     "end-to-end on a candidate sample (≤100) in-browser", dict(size=11.5, bullet="dot")),
    ("submission_metadata.yaml at repo root mirrors the portal metadata",
     dict(size=11.5, bullet="dot")),
])

prs.save("deck_submission.pptx")
print("saved deck_submission.pptx")
